"""Tests for the branded-pptx engine.

No template ships with the Skill, so these tests need a local profile
(`skills/branded-pptx/profiles/default.json`) pointing at one. Without it the
module skips rather than failing — a fresh clone has no corporate template.

Nothing here hard-codes a layout name: layouts are selected by shape, so the
suite exercises the same code paths against whatever template is configured.
"""

import importlib.util
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "skills" / "branded-pptx" / "scripts"

pytest.importorskip("pptx")


def load(name):
    sys.path.insert(0, str(SCRIPTS))
    if name in sys.modules:
        return sys.modules[name]
    spec = importlib.util.spec_from_file_location(name, SCRIPTS / f"{name}.py")
    module = importlib.util.module_from_spec(spec)
    # Register before executing: dataclasses resolve annotations through
    # sys.modules, and the scripts import each other by module name.
    sys.modules[name] = module
    spec.loader.exec_module(module)
    return module


tm = load("textmetrics")
deckprofile = load("deckprofile")
gen = load("generate_deck")
verify = load("verify_pptx")

try:
    PROFILE = deckprofile.resolve()
except deckprofile.ProfileError as exc:
    pytest.skip(f"no template profile configured: {exc}", allow_module_level=True)

COVER = {"slide_type": "cover"}


def build(tmp_path, spec, name="deck.pptx"):
    output = tmp_path / name
    report = gen.generate(spec, output, profile=PROFILE)
    return output, report


def presentation():
    from pptx import Presentation

    return Presentation(str(PROFILE.template))


def layouts():
    return gen.layout_index(presentation())


def body_slots(layout):
    """Text slots a layout offers for content — the title slot excluded, using
    the same resolution rule as the generator."""
    from pptx.util import Inches

    text = [
        s for s in layout.placeholders if s.placeholder_format.type in gen.TEXT_PH_TYPES
    ]
    title = next(
        (s for s in text if s.placeholder_format.idx == PROFILE.title_placeholder_idx), None
    )
    if title is None:
        band = [
            s for s in text if s.top < PROFILE.title_band_bottom and s.width > Inches(4.0)
        ]
        title = min(band, key=lambda s: (s.top, s.left)) if band else None
    return sorted(
        (s for s in text if title is None or s.element is not title.element),
        key=lambda s: (s.top, s.left),
    )


def pick_layout(pictures=0, bodies=0, max_bodies=None, title_idx_is_picture=False):
    """Find a layout in the configured template matching a shape, or skip."""
    from pptx.enum.shapes import PP_PLACEHOLDER

    for name, layout in layouts().items():
        slots = list(layout.placeholders)
        picture_slots = [
            s for s in slots if s.placeholder_format.type == PP_PLACEHOLDER.PICTURE
        ]
        found_bodies = len(body_slots(layout))
        if len(picture_slots) < pictures or found_bodies < bodies:
            continue
        if max_bodies is not None and found_bodies > max_bodies:
            continue
        if title_idx_is_picture and not any(
            s.placeholder_format.idx == PROFILE.title_placeholder_idx for s in picture_slots
        ):
            continue
        return name
    pytest.skip(
        f"template has no layout with >={pictures} picture and >={bodies} body slots"
        + (" where the profile's title index is a picture" if title_idx_is_picture else "")
    )


def png(path, rgb=(0, 96, 240), size=(320, 240)):
    import struct
    import zlib

    width, height = size
    raw = b"".join(b"\x00" + bytes(rgb) * width for _ in range(height))

    def chunk(tag, data):
        body = tag + data
        return struct.pack(">I", len(data)) + body + struct.pack(">I", zlib.crc32(body))

    Path(path).write_bytes(
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0))
        + chunk(b"IDAT", zlib.compress(raw))
        + chunk(b"IEND", b"")
    )
    return str(path)


# --- text metrics -----------------------------------------------------------


def test_cjk_is_twice_as_wide_as_latin():
    assert tm.text_width_pt("中文", 20) == pytest.approx(40.0)
    assert tm.text_width_pt("ab", 20) == pytest.approx(20.8)


def test_wrap_splits_long_cjk_runs():
    lines = tm.wrap_lines("中" * 30, 20, 200.0)
    assert len(lines) == 3
    assert all(tm.text_width_pt(line, 20) <= 200.0 for line in lines)


def test_fit_font_size_returns_none_when_nothing_fits():
    paragraphs = [{"text": "中" * 200, "size_pt": 18}]
    assert tm.fit_font_size(paragraphs, 100.0, 20.0, [18, 16, 14]) is None


# --- profile ----------------------------------------------------------------


def test_missing_profile_names_the_fix(tmp_path):
    with pytest.raises(deckprofile.ProfileError) as excinfo:
        deckprofile.load(tmp_path / "nope.json")
    assert "profile not found" in str(excinfo.value)


def test_profile_pointing_at_a_missing_template_is_rejected(tmp_path):
    path = tmp_path / "p.json"
    path.write_text('{"template": "absent.pptx", "cover_title_shape": "t", "content_layout": "l"}')
    with pytest.raises(deckprofile.ProfileError) as excinfo:
        deckprofile.load(path)
    assert "does not exist" in str(excinfo.value)


def test_unknown_profile_field_is_rejected(tmp_path):
    template = tmp_path / "t.pptx"
    template.write_bytes(PROFILE.template.read_bytes())
    path = tmp_path / "p.json"
    path.write_text(
        '{"template": "t.pptx", "cover_title_shape": "t", "content_layout": "l", "typo": 1}'
    )
    with pytest.raises(deckprofile.ProfileError) as excinfo:
        deckprofile.load(path)
    assert "typo" in str(excinfo.value)


def test_detected_profile_round_trips(tmp_path):
    import inspect_template

    detected = inspect_template.detect_profile(presentation(), PROFILE.template)
    assert detected["cover_title_shape"] == PROFILE.cover_title_shape
    assert detected["content_layout"] == PROFILE.content_layout
    assert detected["title_placeholder_idx"] == PROFILE.title_placeholder_idx


# --- spec validation --------------------------------------------------------


@pytest.mark.parametrize(
    "page, message",
    [
        ({"slide_type": "content", "title": "T", "content_type": "table"}, "table.headers"),
        (
            {
                "slide_type": "content",
                "title": "T",
                "content_type": "table",
                "table": {"headers": ["a", "b"], "rows": [["x"]]},
            },
            "headers declare 2",
        ),
        ({"slide_type": "content", "title": "T", "content_type": "nope"}, "content_type must be"),
        ({"slide_type": "content", "content": ["x"]}, "title is required"),
        (
            {"slide_type": "content", "title": "T", "content_type": "layout", "layout": "不存在"},
            "not in the template",
        ),
    ],
)
def test_invalid_specs_raise_instead_of_producing_blank_slides(tmp_path, page, message):
    with pytest.raises(gen.SpecError) as excinfo:
        build(tmp_path, {"title": "T", "outline": [COVER, page]})
    assert message in str(excinfo.value)


def test_picture_only_layout_without_images_is_rejected(tmp_path):
    layout = pick_layout(pictures=1, max_bodies=0)
    spec = {
        "title": "T",
        "outline": [
            {"slide_type": "content", "title": "图页", "content_type": "layout", "layout": layout}
        ],
    }
    with pytest.raises(gen.SpecError) as excinfo:
        build(tmp_path, spec)
    assert "empty page" in str(excinfo.value)


# --- cover ------------------------------------------------------------------


def _cover_title_shape(path):
    from pptx import Presentation

    slide = Presentation(str(path)).slides[0]
    return next(s for s in slide.shapes if s.name == PROFILE.cover_title_shape)


def test_cover_lines_share_one_size_and_one_paragraph_style(tmp_path):
    output, _ = build(tmp_path, {"title": "第一行标题\n第二行副标题更长一些", "outline": [COVER]})
    shape = _cover_title_shape(output)
    sizes = {
        run.font.size.pt
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    }
    spacings = {paragraph.line_spacing for paragraph in shape.text_frame.paragraphs}
    assert len(sizes) == 1
    assert len(spacings) == 1


def test_long_cover_title_wraps_instead_of_running_off_the_slide(tmp_path):
    title = "某某某平台在某某行业的规模化落地实践与下一阶段路线图说明"
    output, report = build(tmp_path, {"title": title, "outline": [COVER]})
    shape = _cover_title_shape(output)
    body_pr = shape.text_frame._txBody.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr"
    )
    assert body_pr.get("wrap") == "square"
    assert len(shape.text_frame.paragraphs) > 1
    assert any("wrap mode" in note for note in report["notes"])
    assert not [f for f in verify.check_file(output) if f["level"] == "error"]


def test_date_field_is_replaced_with_static_text(tmp_path):
    if not PROFILE.cover_date_shape:
        pytest.skip("profile has no cover date shape")
    output, _ = build(tmp_path, {"title": "T", "date": "2026年8月5日", "outline": [COVER]})
    from pptx import Presentation

    slide = Presentation(str(output)).slides[0]
    shape = next(s for s in slide.shapes if s.name == PROFILE.cover_date_shape)
    assert shape.text_frame.text == "2026年8月5日"
    assert shape.text_frame._txBody.findall(
        ".//{http://schemas.openxmlformats.org/drawingml/2006/main}fld"
    ) == []


# --- pagination -------------------------------------------------------------


def test_oversized_table_paginates_and_stays_on_the_slide(tmp_path):
    rows = [[f"客户{i}", "场景", "效果"] for i in range(40)]
    spec = {
        "title": "T",
        "outline": [
            {
                "slide_type": "content",
                "title": "清单",
                "content_type": "table",
                "table": {"headers": ["客户", "场景", "效果"], "rows": rows},
            }
        ],
    }
    output, report = build(tmp_path, spec)
    assert report["slides"] > 1
    assert not [f for f in verify.check_file(output) if f["level"] == "error"]

    from pptx import Presentation

    deck = Presentation(str(output))
    for slide in deck.slides:
        table = next(s for s in slide.shapes if s.has_table)
        bottom = table.top + sum(row.height for row in table.table.rows)
        assert bottom <= deck.slide_height


def test_many_bullets_paginate_with_a_continuation_title(tmp_path):
    items = [f"第 {i} 条要点，" + "内容" * 20 for i in range(30)]
    output, report = build(
        tmp_path,
        {"title": "T", "outline": [{"slide_type": "content", "title": "要点", "content": items}]},
    )
    assert report["pages_added"] >= 1
    assert not [f for f in verify.check_file(output) if f["level"] == "error"]

    from pptx import Presentation

    titles = [
        gen.find_title_placeholder(slide).text_frame.text
        for slide in Presentation(str(output)).slides
    ]
    assert titles[0] == "要点"
    assert all(title.endswith(gen.CONTINUATION_SUFFIX) for title in titles[1:])


# --- layouts ----------------------------------------------------------------


def test_title_lands_in_a_text_slot_even_when_the_profile_index_is_a_picture(tmp_path):
    """Templates reuse placeholder indexes across types; resolving by index
    alone writes the page title into an image box."""
    layout = pick_layout(pictures=1, bodies=1, title_idx_is_picture=True)
    spec = {
        "title": "T",
        "outline": [
            {
                "slide_type": "content",
                "title": "版式页",
                "content_type": "layout",
                "layout": layout,
                "content": [f"文案{i}" for i in range(8)],
            }
        ],
    }
    output, _ = build(tmp_path, spec)
    from pptx import Presentation

    slide = Presentation(str(output)).slides[0]
    title = gen.find_title_placeholder(slide)
    assert title.placeholder_format.type in gen.TEXT_PH_TYPES
    assert title.text_frame.text == "版式页"
    assert not [f for f in verify.check_file(output) if f["level"] == "error"]


def test_slots_address_placeholders_by_index(tmp_path):
    layout = pick_layout(bodies=2)
    slots = body_slots(layouts()[layout])
    first, second = slots[0].placeholder_format.idx, slots[1].placeholder_format.idx
    spec = {
        "title": "T",
        "outline": [
            {
                "slide_type": "content",
                "title": "版式页",
                "content_type": "layout",
                "layout": layout,
                "slots": {str(first): "甲", str(second): "乙"},
                "images": [png(tmp_path / "a.png")],
            }
        ],
    }
    output, _ = build(tmp_path, spec)
    from pptx import Presentation

    slide = Presentation(str(output)).slides[0]
    filled = {
        s.placeholder_format.idx: s.text_frame.text
        for s in slide.placeholders
        if s.has_text_frame and s.text_frame.text
    }
    assert filled[first] == "甲"
    assert filled[second] == "乙"


def test_unknown_slot_index_is_rejected(tmp_path):
    layout = pick_layout(bodies=1)
    spec = {
        "title": "T",
        "outline": [
            {
                "slide_type": "content",
                "title": "版式页",
                "content_type": "layout",
                "layout": layout,
                "slots": {"99": "x"},
            }
        ],
    }
    with pytest.raises(gen.SpecError) as excinfo:
        build(tmp_path, spec)
    assert "99" in str(excinfo.value)


def test_pictures_fill_picture_slots(tmp_path):
    layout = pick_layout(pictures=2)
    images = [png(tmp_path / "a.png"), png(tmp_path / "b.png", rgb=(200, 40, 40))]
    spec = {
        "title": "T",
        "outline": [
            {
                "slide_type": "content",
                "title": "图页",
                "content_type": "layout",
                "layout": layout,
                "images": images,
            }
        ],
    }
    output, _ = build(tmp_path, spec)
    from pptx import Presentation

    slide = Presentation(str(output)).slides[0]
    assert sum(1 for s in slide.shapes if s.element.tag.endswith("}pic")) == 2
    assert not [f for f in verify.check_file(output) if f["level"] == "error"]


# --- deck assembly ----------------------------------------------------------


def test_only_requested_slides_survive_and_order_matches_the_outline(tmp_path):
    spec = {
        "title": "T",
        "outline": [
            COVER,
            {"slide_type": "content", "title": "第一页", "content": ["a"]},
            {"slide_type": "content", "title": "第二页", "content": ["b"]},
        ],
    }
    output, report = build(tmp_path, spec)
    assert report["slides"] == 3

    from pptx import Presentation

    slides = list(Presentation(str(output)).slides)
    assert len(slides) == 3
    assert [gen.find_title_placeholder(s).text_frame.text for s in slides[1:]] == [
        "第一页",
        "第二页",
    ]


def test_repeated_cover_is_cloned_with_its_images(tmp_path):
    from pptx import Presentation

    source_pictures = sum(
        1
        for s in presentation().slides[PROFILE.cover_slide_index].shapes
        if s.element.tag.endswith("}pic")
    )
    output, _ = build(tmp_path, {"title": "T", "outline": [COVER, COVER]})
    counts = [
        sum(1 for s in slide.shapes if s.element.tag.endswith("}pic"))
        for slide in Presentation(str(output)).slides
    ]
    assert counts == [source_pictures, source_pictures]


def test_example_spec_generates_cleanly(tmp_path):
    output, _ = build(tmp_path, gen.EXAMPLE_SPEC)
    assert not [f for f in verify.check_file(output) if f["level"] == "error"]


# --- verifier ---------------------------------------------------------------


def test_verifier_flags_text_that_overflows_its_box(tmp_path):
    from pptx import Presentation
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.util import Inches

    output, _ = build(tmp_path, gen.EXAMPLE_SPEC)
    deck = Presentation(str(output))
    box = deck.slides[1].shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(3), Inches(0.4))
    box.text_frame.word_wrap = True
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    box.text_frame.paragraphs[0].add_run().text = "超长文本" * 40
    broken = tmp_path / "broken.pptx"
    deck.save(str(broken))

    assert any(
        "overflow" in f["message"] for f in verify.check_file(broken) if f["level"] == "error"
    )


def test_verifier_flags_text_in_a_picture_placeholder(tmp_path):
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER

    layout = pick_layout(pictures=1)
    deck = presentation()
    slide = deck.slides.add_slide(gen.layout_index(deck)[layout])
    picture_slot = next(
        s
        for s in slide.placeholders
        if s.placeholder_format.type == PP_PLACEHOLDER.PICTURE
    )
    picture_slot.text_frame.paragraphs[0].add_run().text = "误写入图片占位符"
    broken = tmp_path / "broken.pptx"
    deck.save(str(broken))

    assert any(
        "picture/chart placeholder" in f["message"] for f in verify.check_file(broken)
    )
