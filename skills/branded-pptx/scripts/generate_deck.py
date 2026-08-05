#!/usr/bin/env python3
"""Generate a deck on your own corporate template from a JSON spec.

Slide assembly happens in this process, so the Skill depends on nothing but
python-pptx — no other Skill, no client-specific home directory. The template
and the handful of facts about it live outside the code, in a profile
(see deckprofile.py); the rendering engine itself is template-agnostic.

Design rules enforced by this script (not merely documented):
  - the cover title uses one font size for every line, one paragraph style for
    every line, and shrinks until the longest line fits the profile's safe width;
  - body text and tables are measured before they are written, and content that
    cannot fit is paginated onto continuation slides rather than silently
    overflowing;
  - brand colours and fonts come from the template theme, never from hard-coded
    RGB values or font names;
  - a spec that cannot be rendered faithfully raises SpecError instead of
    producing a blank slide.

Usage:
    python3 generate_deck.py spec.json --profile profiles/default.json
    python3 generate_deck.py spec.json --output deck.pptx
    python3 generate_deck.py --example > spec.json
"""

from __future__ import annotations

import argparse
import json
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
import deckprofile  # noqa: E402
import textmetrics as tm  # noqa: E402

SCRIPT_DIR = Path(__file__).resolve().parent
SKILL_DIR = SCRIPT_DIR.parent

# The template facts for the deck currently being built. generate() sets this
# from the resolved profile before any rendering happens.
PROFILE = deckprofile.Profile(template=Path("unset"))

# Measured once per deck, from the presentation and its content layout, so that
# pagination can be planned before any slide exists.
_SLIDE_WIDTH: int = int(Inches(13.333))
_CONTENT_TOP: int = int(Inches(1.1))

TEXT_PH_TYPES = {
    PP_PLACEHOLDER.BODY,
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
    PP_PLACEHOLDER.SUBTITLE,
    PP_PLACEHOLDER.OBJECT,
}

# A cover title smaller than 24pt looks weak; below that the box wraps instead.
COVER_SIZES = [36, 34, 32, 30, 28, 26, 24]

# Largest size that fits wins, so a sparse page gets bigger, more readable type
# instead of a small block floating in whitespace.
BULLET_SIZES = [22, 20, 18, 17, 16, 15]
COLUMN_BODY_SIZES = [18, 17, 16, 15, 14, 13]
MIXED_BULLET_SIZES = [18, 17, 16, 15, 14]
TABLE_HEADER_SIZE = 15
TABLE_BODY_SIZE = 14
TABLE_HEADER_HEIGHT = Inches(0.48)
TABLE_ROW_HEIGHT = Inches(0.42)
CONTINUATION_SUFFIX = "（续）"

VALID_CONTENT_TYPES = {"bullets", "table", "columns", "mixed", "layout"}


class SpecError(ValueError):
    """The spec cannot be rendered faithfully; the caller must fix the spec."""


# ----------------------------------------------------------------------------
# Spec validation
# ----------------------------------------------------------------------------


def validate_spec(spec: dict, template_layouts: set[str] | None = None) -> None:
    if not isinstance(spec, dict):
        raise SpecError("spec must be a JSON object")
    if not spec.get("title"):
        raise SpecError("spec.title is required (it is the cover title)")
    outline = spec.get("outline")
    if not isinstance(outline, list) or not outline:
        raise SpecError("spec.outline must be a non-empty list of pages")

    for index, page in enumerate(outline):
        where = f"outline[{index}]"
        if not isinstance(page, dict):
            raise SpecError(f"{where} must be an object")
        slide_type = page.get("slide_type", "content")
        if slide_type not in ("cover", "content"):
            raise SpecError(f"{where}.slide_type must be 'cover' or 'content'")
        if slide_type == "cover":
            continue

        if not page.get("title"):
            raise SpecError(f"{where}.title is required on content pages")
        content_type = page.get("content_type", "bullets")
        if content_type not in VALID_CONTENT_TYPES:
            raise SpecError(
                f"{where}.content_type must be one of "
                f"{', '.join(sorted(VALID_CONTENT_TYPES))}; got {content_type!r}"
            )

        if content_type in ("table", "mixed"):
            table = page.get("table")
            if not isinstance(table, dict) or not table.get("headers"):
                raise SpecError(
                    f"{where} uses content_type={content_type!r} but has no "
                    "table.headers; add the table or change content_type"
                )
            rows = table.get("rows")
            if not isinstance(rows, list) or not rows:
                raise SpecError(f"{where}.table.rows must be a non-empty list")
            width = len(table["headers"])
            for row_index, row in enumerate(rows):
                if not isinstance(row, list):
                    raise SpecError(f"{where}.table.rows[{row_index}] must be a list")
                if len(row) != width:
                    raise SpecError(
                        f"{where}.table.rows[{row_index}] has {len(row)} cells but "
                        f"headers declare {width}"
                    )

        if content_type in ("bullets", "columns", "mixed") and not page.get("content"):
            raise SpecError(
                f"{where} uses content_type={content_type!r} but content is empty"
            )

        if content_type == "layout":
            layout = page.get("layout")
            if not layout:
                raise SpecError(
                    f"{where} uses content_type='layout' but no layout name is set; "
                    "see references/layouts.md"
                )
            if template_layouts is not None and layout not in template_layouts:
                raise SpecError(
                    f"{where}.layout={layout!r} is not in the template; "
                    "run scripts/inspect_template.py --layouts to list valid names"
                )


# ----------------------------------------------------------------------------
# Slide plumbing
# ----------------------------------------------------------------------------


def layout_index(prs: Presentation) -> dict[str, object]:
    """Map layout name -> layout, preferring the master used by the sample slides."""
    index: dict[str, object] = {}
    preferred_master = prs.slides[1].slide_layout.slide_master if len(prs.slides) > 1 else None
    masters = list(prs.slide_masters)
    if preferred_master is not None:
        masters.sort(key=lambda master: master is not preferred_master)
    for master in masters:
        for layout in master.slide_layouts:
            index.setdefault(layout.name, layout)
    return index


def clone_slide(prs: Presentation, source):
    """Append a copy of `source`, remapping relationship ids so images survive."""
    new_slide = prs.slides.add_slide(source.slide_layout)
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    rel_map: dict[str, str] = {}
    for r_id, rel in source.part.rels.items():
        if rel.is_external:
            rel_map[r_id] = new_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            rel_map[r_id] = new_slide.part.relate_to(rel.target_part, rel.reltype)

    r_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    for shape in source.shapes:
        element = deepcopy(shape._element)
        for node in element.iter():
            for attr, value in list(node.attrib.items()):
                if attr.startswith(r_ns) and value in rel_map:
                    node.set(attr, rel_map[value])
        new_slide.shapes._spTree.append(element)
    return new_slide


def remove_slides(prs: Presentation, slides: list) -> None:
    """Remove slides by identity, dropping both the sldId entry and the relationship."""
    targets = {id(slide.part) for slide in slides}
    slide_id_list = prs.slides._sldIdLst
    for slide_id in list(slide_id_list):
        r_id = slide_id.get(qn("r:id"))
        part = prs.part.rels[r_id].target_part
        if id(part) in targets:
            slide_id_list.remove(slide_id)
            prs.part.drop_rel(r_id)


def reorder_slides(prs: Presentation, ordered_slides: list) -> None:
    slide_id_list = prs.slides._sldIdLst
    entries = {}
    for slide_id in list(slide_id_list):
        r_id = slide_id.get(qn("r:id"))
        entries[id(prs.part.rels[r_id].target_part)] = slide_id
        slide_id_list.remove(slide_id)
    for slide in ordered_slides:
        slide_id_list.append(entries[id(slide.part)])


def find_shape_by_name(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def find_title_placeholder(slide):
    """Resolve the page-title placeholder.

    The profile records the usual index, but templates are inconsistent: picture
    layouts often renumber the title slot, and a layout may use the profile's
    index for a *picture* instead — so the placeholder type is checked before
    the index. Fall back to the topmost wide text slot in the title band.
    """
    text_placeholders = [
        shape
        for shape in slide.shapes
        if shape.is_placeholder and shape.placeholder_format.type in TEXT_PH_TYPES
    ]
    for shape in text_placeholders:
        if shape.placeholder_format.idx == PROFILE.title_placeholder_idx:
            return shape
    band = [
        shape
        for shape in text_placeholders
        if shape.top < PROFILE.title_band_bottom and shape.width > Inches(4.0)
    ]
    return min(band, key=lambda shape: (shape.top, shape.left), default=None)


def content_area(slide) -> tuple[int, int, int, int]:
    """(left, top, width, height) of the usable body area, in EMU."""
    title = find_title_placeholder(slide)
    top = (
        title.top + title.height + PROFILE.content_gap_below_title
        if title
        else PROFILE.content_gap_below_title
    )
    left = PROFILE.content_margin_x
    width = _SLIDE_WIDTH - 2 * left
    height = PROFILE.content_bottom - top
    return int(left), int(top), int(width), int(height)


# ----------------------------------------------------------------------------
# Cover
# ----------------------------------------------------------------------------


def cover_title_budget(shape) -> float:
    """Width available to the cover title before it collides with the photo, in points."""
    return Emu(PROFILE.cover_safe_right - PROFILE.cover_safe_margin - shape.left).pt


def fit_cover_size(lines: list[str], budget_pt: float) -> tuple[int, bool]:
    """Pick the largest size where every line fits. Returns (size, needs_wrap)."""
    for size in COVER_SIZES:
        if all(tm.text_width_pt(line, size) <= budget_pt for line in lines):
            return size, False
    return COVER_SIZES[-1], True


def set_cover_title(shape, title_text: str) -> dict:
    """Rewrite the cover title with one font size and one paragraph style per line.

    The template run carries `+mj-ea` and a theme colour; both are preserved by
    cloning the original run properties and overriding only the size.
    """
    text_frame = shape.text_frame
    body = text_frame._txBody
    paragraphs = body.findall(qn("a:p"))
    if not paragraphs:
        raise SpecError("cover title shape has no paragraph to clone formatting from")

    proto_p = paragraphs[0]
    proto_pPr = proto_p.find(qn("a:pPr"))
    proto_run = proto_p.find(qn("a:r"))
    proto_rPr = proto_run.find(qn("a:rPr")) if proto_run is not None else None

    lines = [line.strip() for line in title_text.split("\n") if line.strip()]
    if not lines:
        raise SpecError("cover title is empty")

    budget_pt = cover_title_budget(shape)
    size_pt, needs_wrap = fit_cover_size(lines, budget_pt)

    if needs_wrap:
        # Longest line still does not fit at the minimum size: switch the box
        # from grow-sideways to wrap-inside-a-fixed-width so nothing runs off.
        body_pr = body.find(qn("a:bodyPr"))
        body_pr.set("wrap", "square")
        shape.width = PROFILE.cover_safe_right - PROFILE.cover_safe_margin - shape.left
        wrapped: list[str] = []
        for line in lines:
            wrapped.extend(tm.wrap_lines(line, size_pt, budget_pt))
        lines = wrapped

    for paragraph in paragraphs:
        body.remove(paragraph)

    for line in lines:
        p = body.makeelement(qn("a:p"), {})
        if proto_pPr is not None:
            p.append(deepcopy(proto_pPr))
        run = p.makeelement(qn("a:r"), {})
        if proto_rPr is not None:
            run_props = deepcopy(proto_rPr)
        else:
            run_props = run.makeelement(qn("a:rPr"), {"lang": "zh-CN"})
        run_props.set("sz", str(int(size_pt * 100)))
        run.append(run_props)
        text_node = run.makeelement(qn("a:t"), {})
        text_node.text = line
        run.append(text_node)
        p.append(run)
        body.append(p)

    return {"lines": len(lines), "size_pt": size_pt, "wrapped": needs_wrap}


def set_cover_date(shape, date_text: str) -> None:
    """Replace the auto-updating date field with static text, keeping its style."""
    paragraph = shape.text_frame.paragraphs[0]._p
    proto_rPr = None
    for child in list(paragraph):
        tag = child.tag.split("}")[-1]
        if tag in ("fld", "r"):
            if proto_rPr is None:
                found = child.find(qn("a:rPr"))
                if found is not None:
                    proto_rPr = deepcopy(found)
            paragraph.remove(child)

    run = paragraph.makeelement(qn("a:r"), {})
    if proto_rPr is not None:
        for attr in ("smtClean", "dirty"):
            proto_rPr.attrib.pop(attr, None)
        run.append(proto_rPr)
    text_node = run.makeelement(qn("a:t"), {})
    text_node.text = date_text
    run.append(text_node)
    paragraph.append(run)


# ----------------------------------------------------------------------------
# Content renderers
# ----------------------------------------------------------------------------


def normalize_item(item) -> dict:
    if isinstance(item, dict):
        return {
            "text": str(item.get("text", "")),
            "level": int(item.get("level", 0)),
            "bold": bool(item.get("bold", False)),
        }
    return {"text": str(item), "level": 0, "bold": False}


def bullet_paragraph_specs(items: list[dict], base_size: int) -> list[dict]:
    return [
        {
            "text": item["text"],
            "size_pt": base_size if item["level"] == 0 else base_size - 2,
            "space_before_pt": 8 if item["level"] == 0 else 4,
            "space_after_pt": 4,
            "line_spacing": 1.15,
        }
        for item in items
    ]


def set_title(slide, text: str) -> None:
    placeholder = find_title_placeholder(slide)
    if placeholder is None:
        raise SpecError(
            f"layout {slide.slide_layout.name!r} has no title slot at idx="
            f"{PROFILE.title_placeholder_idx} and none in the title band; "
            "re-derive the profile with scripts/inspect_template.py --profile"
        )
    text_frame = placeholder.text_frame
    for extra in text_frame.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    paragraph = text_frame.paragraphs[0]
    for run in list(paragraph.runs)[1:]:
        run._r.getparent().remove(run._r)
    if paragraph.runs:
        paragraph.runs[0].text = text
    else:
        paragraph.add_run().text = text


def render_bullets(slide, items: list[dict], area=None, sizes=None) -> None:
    left, top, width, height = area or content_area(slide)
    sizes = sizes or BULLET_SIZES
    specs = bullet_paragraph_specs(items, sizes[0])
    size = tm.fit_font_size(specs, Emu(width).pt, Emu(height).pt, sizes) or sizes[-1]

    box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE

    for index, item in enumerate(items):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        run = paragraph.add_run()
        run.text = item["text"]
        run.font.size = Pt(size if item["level"] == 0 else size - 2)
        run.font.bold = item["bold"]
        paragraph.level = item["level"]
        paragraph.space_before = Pt(8 if item["level"] == 0 else 4)
        paragraph.space_after = Pt(4)


def render_table(slide, table_data: dict, area=None) -> None:
    left, top, width, height = area or content_area(slide)
    headers = table_data["headers"]
    rows = table_data["rows"]

    shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, TABLE_HEADER_HEIGHT)
    table = shape.table
    column_width = int(width / len(headers))
    for index in range(len(headers)):
        table.columns[index].width = column_width

    for column, header in enumerate(headers):
        cell = table.cell(0, column)
        _fill_cell(cell, header, TABLE_HEADER_SIZE, bold=True, align=PP_ALIGN.CENTER)
        cell.fill.solid()
        cell.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        cell.text_frame.paragraphs[0].runs[0].font.color.theme_color = MSO_THEME_COLOR.BACKGROUND_1

    for row_index, row in enumerate(rows):
        for column, value in enumerate(row):
            _fill_cell(table.cell(row_index + 1, column), str(value), TABLE_BODY_SIZE)

    table.rows[0].height = TABLE_HEADER_HEIGHT
    for row_index in range(1, len(rows) + 1):
        table.rows[row_index].height = TABLE_ROW_HEIGHT


def _fill_cell(cell, text: str, size_pt: int, bold: bool = False, align=None) -> None:
    cell.text = ""
    paragraph = cell.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if align is not None:
        paragraph.alignment = align
    cell.text_frame.word_wrap = True


def render_columns(slide, items: list, column_count: int) -> None:
    left, top, width, height = content_area(slide)
    structured = bool(items) and isinstance(items[0], dict) and "items" in items[0]
    if structured:
        columns = items
    else:
        per_column = max(1, (len(items) + column_count - 1) // column_count)
        columns = [
            {"items": items[index : index + per_column]}
            for index in range(0, len(items), per_column)
        ]
    count = len(columns)
    gap = Inches(0.3)
    column_width = int((width - gap * (count - 1)) / count)

    for index, column in enumerate(columns):
        column_left = left + index * (column_width + gap)
        box = slide.shapes.add_textbox(column_left, top, column_width, height)
        text_frame = box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE

        entries = [normalize_item(entry) for entry in column.get("items", [])]
        specs = []
        if column.get("title"):
            specs.append({"text": column["title"], "size_pt": COLUMN_BODY_SIZES[0] + 3, "space_after_pt": 10})
        specs.extend(
            {"text": entry["text"], "size_pt": COLUMN_BODY_SIZES[0], "space_before_pt": 6, "space_after_pt": 3}
            for entry in entries
        )
        size = tm.fit_font_size(specs, Emu(column_width).pt, Emu(height).pt, COLUMN_BODY_SIZES)
        size = size or COLUMN_BODY_SIZES[-1]

        first = True
        if column.get("title"):
            paragraph = text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = column["title"]
            run.font.size = Pt(size + 3)
            run.font.bold = True
            run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
            paragraph.space_after = Pt(10)
            first = False

        for entry in entries:
            paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
            first = False
            run = paragraph.add_run()
            run.text = entry["text"]
            run.font.size = Pt(size)
            run.font.bold = entry["bold"]
            paragraph.space_before = Pt(6)
            paragraph.space_after = Pt(3)


def render_mixed(slide, items: list[dict], table_data: dict) -> None:
    left, top, width, height = content_area(slide)
    specs = bullet_paragraph_specs(items, MIXED_BULLET_SIZES[0])
    bullet_height = int(min(height * 0.4, Pt(tm.block_height_pt(specs, Emu(width).pt)) + Pt(12)))
    render_bullets(slide, items, area=(left, top, width, bullet_height), sizes=MIXED_BULLET_SIZES)
    table_top = top + bullet_height + Inches(0.2)
    render_table(slide, table_data, area=(left, table_top, width, top + height - table_top))


def render_layout_page(slide, page: dict) -> list[str]:
    """Fill a template layout: title, then body placeholders, then pictures.

    Body placeholders are filled in reading order (top to bottom, left to
    right). A `slots` map keyed by placeholder idx overrides that order when the
    layout's reading order is not what the content needs.
    """
    title_shape = find_title_placeholder(slide)
    set_title(slide, page["title"])
    entries = [normalize_item(entry) for entry in page.get("content", [])]
    slots = {str(key): value for key, value in (page.get("slots") or {}).items()}
    images = list(page.get("images", []))
    notes: list[str] = []

    bodies, pictures, charts = [], [], []
    for shape in slide.placeholders:
        if title_shape is not None and shape.element is title_shape.element:
            continue
        kind = shape.placeholder_format.type
        if kind == PP_PLACEHOLDER.PICTURE:
            pictures.append(shape)
        elif kind == PP_PLACEHOLDER.CHART:
            charts.append(shape)
        else:
            bodies.append(shape)

    def reading_order(shape):
        return (round(shape.top / 100000), shape.left)

    bodies.sort(key=reading_order)
    pictures.sort(key=reading_order)

    if slots:
        unknown = slots.keys() - {str(shape.placeholder_format.idx) for shape in bodies}
        if unknown:
            raise SpecError(
                f"slots reference placeholder idx {sorted(unknown)} which layout "
                f"{page['layout']!r} does not have"
            )
        filled = []
        for shape in bodies:
            key = str(shape.placeholder_format.idx)
            if key in slots:
                _fill_placeholder(shape, normalize_item(slots[key]))
                filled.append(shape)
        bodies = [shape for shape in bodies if shape not in filled]

    for shape, entry in zip(bodies, entries):
        _fill_placeholder(shape, entry)

    filled_bodies = len(slots) + min(len(bodies), len(entries))
    if not filled_bodies and not images:
        raise SpecError(
            f"layout {page['layout']!r} would render as a title on an empty page: it has "
            f"{len(bodies) + len(slots)} text slot(s) and {len(pictures)} picture slot(s), "
            "and neither content nor images reached them — supply page.images, use "
            "page.slots, or pick another layout"
        )
    if pictures and not images:
        notes.append(
            f"layout {page['layout']!r} has {len(pictures)} picture slot(s) but no images "
            "were supplied; the empty slots were removed"
        )

    for shape, image in zip(pictures + charts, images):
        path = Path(image)
        if not path.exists():
            raise SpecError(f"image not found: {image}")
        if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            shape.insert_picture(str(path))
        else:
            box = (shape.left, shape.top, shape.width, shape.height)
            shape._element.getparent().remove(shape._element)
            slide.shapes.add_picture(str(path), *box)

    # Unused placeholders would render as "click to add" prompts: drop them.
    for shape in bodies[len(entries) :] + (pictures + charts)[len(images) :]:
        if shape._element.getparent() is not None:
            shape._element.getparent().remove(shape._element)
    return notes


def _fill_placeholder(shape, entry: dict) -> None:
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    run = text_frame.paragraphs[0].add_run()
    run.text = entry["text"]
    if entry["bold"]:
        run.font.bold = True


# ----------------------------------------------------------------------------
# Pagination
# ----------------------------------------------------------------------------


def _measure_content_top(layout) -> int:
    """Top of the body area, from the content layout's own title placeholder."""
    candidates = [
        shape
        for shape in layout.placeholders
        if shape.placeholder_format.type in TEXT_PH_TYPES
        and shape.top < PROFILE.title_band_bottom
        and shape.width > Inches(4.0)
    ]
    if not candidates:
        return int(PROFILE.content_gap_below_title)
    title = min(candidates, key=lambda shape: (shape.top, shape.left))
    return int(title.top + title.height + PROFILE.content_gap_below_title)


def _default_content_metrics() -> tuple[float, float]:
    """Body width/height of the default content layout, in points."""
    width = _SLIDE_WIDTH - 2 * PROFILE.content_margin_x
    return Emu(int(width)).pt, Emu(int(PROFILE.content_bottom - _CONTENT_TOP)).pt


def paginate_bullets(items: list[dict]) -> list[list[dict]]:
    width_pt, height_pt = _default_content_metrics()
    pages, current = [], []
    for item in items:
        candidate = current + [item]
        specs = bullet_paragraph_specs(candidate, BULLET_SIZES[-1])
        if current and tm.block_height_pt(specs, width_pt) > height_pt:
            pages.append(current)
            current = [item]
        else:
            current = candidate
    if current:
        pages.append(current)
    return pages


def rows_per_page(reserved_pt: float = 0.0) -> int:
    _, height_pt = _default_content_metrics()
    usable = height_pt - reserved_pt - Emu(int(TABLE_HEADER_HEIGHT)).pt
    return max(1, int(usable // Emu(int(TABLE_ROW_HEIGHT)).pt))


def paginate_table(table_data: dict, reserved_pt: float = 0.0) -> list[dict]:
    limit = rows_per_page(reserved_pt)
    rows = table_data["rows"]
    return [
        {"headers": table_data["headers"], "rows": rows[index : index + limit]}
        for index in range(0, len(rows), limit)
    ]


def expand_outline(outline: list[dict]) -> list[dict]:
    """Split oversized pages into continuation pages before any slide is created."""
    pages: list[dict] = []
    for page in outline:
        if page.get("slide_type") == "cover":
            pages.append(dict(page))
            continue

        content_type = page.get("content_type", "bullets")
        title = page["title"]

        if content_type == "bullets":
            items = [normalize_item(entry) for entry in page.get("content", [])]
            chunks = paginate_bullets(items)
            for index, chunk in enumerate(chunks):
                pages.append(
                    {
                        **page,
                        "title": title if index == 0 else title + CONTINUATION_SUFFIX,
                        "content": chunk,
                    }
                )
        elif content_type == "table":
            for index, part in enumerate(paginate_table(page["table"])):
                pages.append(
                    {
                        **page,
                        "title": title if index == 0 else title + CONTINUATION_SUFFIX,
                        "table": part,
                    }
                )
        elif content_type == "mixed":
            items = [normalize_item(entry) for entry in page.get("content", [])]
            _, height_pt = _default_content_metrics()
            specs = bullet_paragraph_specs(items, MIXED_BULLET_SIZES[0])
            width_pt, _ = _default_content_metrics()
            reserved = min(height_pt * 0.4, tm.block_height_pt(specs, width_pt) + 12) + 14
            parts = paginate_table(page["table"], reserved_pt=reserved)
            for index, part in enumerate(parts):
                pages.append(
                    {
                        **page,
                        "title": title if index == 0 else title + CONTINUATION_SUFFIX,
                        "content": items if index == 0 else [],
                        "content_type": "mixed" if index == 0 else "table",
                        "table": part,
                    }
                )
        else:
            pages.append(dict(page))
    return pages


# ----------------------------------------------------------------------------
# Orchestration
# ----------------------------------------------------------------------------


def generate(
    spec: dict,
    output_path: str | Path,
    profile: "deckprofile.Profile | str | Path | None" = None,
) -> dict:
    global PROFILE, _SLIDE_WIDTH, _CONTENT_TOP

    if not isinstance(profile, deckprofile.Profile):
        profile = deckprofile.resolve(profile, spec)
    PROFILE = profile

    prs = Presentation(str(profile.template))
    _SLIDE_WIDTH = prs.slide_width
    layouts = layout_index(prs)
    validate_spec(spec, set(layouts))

    if profile.content_layout not in layouts:
        raise SpecError(
            f"profile content_layout={profile.content_layout!r} is not in "
            f"{profile.template.name}; re-derive the profile with "
            "scripts/inspect_template.py --profile"
        )
    _CONTENT_TOP = _measure_content_top(layouts[profile.content_layout])

    sample_slides = list(prs.slides)
    if profile.cover_slide_index >= len(sample_slides):
        raise SpecError(
            f"profile cover_slide_index={profile.cover_slide_index} but "
            f"{profile.template.name} has {len(sample_slides)} slide(s)"
        )
    cover_source = sample_slides[profile.cover_slide_index]
    pages = expand_outline(spec["outline"])

    built: list = []
    cover_used = False
    for page in pages:
        if page.get("slide_type") == "cover":
            slide = cover_source if not cover_used else clone_slide(prs, cover_source)
            cover_used = True
        else:
            name = (
                page.get("layout")
                if page.get("content_type") == "layout"
                else profile.content_layout
            )
            slide = prs.slides.add_slide(layouts[name])
        built.append(slide)

    kept = {id(slide.part) for slide in built}
    unused = [slide for slide in sample_slides if id(slide.part) not in kept]
    remove_slides(prs, unused)
    reorder_slides(prs, built)

    report = {"slides": len(built), "pages_added": len(pages) - len(spec["outline"]), "notes": []}
    date_text = spec.get("date")

    for page, slide in zip(pages, built):
        if page.get("slide_type") == "cover":
            title_shape = find_shape_by_name(slide, profile.cover_title_shape)
            if title_shape is None:
                raise SpecError(
                    f"cover shape {profile.cover_title_shape!r} is missing from "
                    f"{profile.template.name}; re-derive the profile with "
                    "scripts/inspect_template.py --profile"
                )
            info = set_cover_title(title_shape, spec["title"])
            if info["wrapped"]:
                report["notes"].append(
                    "cover title exceeded the panel width; the box was switched to "
                    "wrap mode — consider a shorter title"
                )
            if date_text and profile.cover_date_shape:
                date_shape = find_shape_by_name(slide, profile.cover_date_shape)
                if date_shape is None:
                    raise SpecError(
                        f"cover shape {profile.cover_date_shape!r} is missing from "
                        f"{profile.template.name}"
                    )
                set_cover_date(date_shape, date_text)
            continue

        content_type = page.get("content_type", "bullets")
        if content_type == "layout":
            report["notes"].extend(render_layout_page(slide, page))
            continue

        set_title(slide, page["title"])
        if content_type == "table":
            render_table(slide, page["table"])
        elif content_type == "columns":
            render_columns(slide, page.get("content", []), int(page.get("columns", 2)))
        elif content_type == "mixed":
            render_mixed(slide, [normalize_item(i) for i in page.get("content", [])], page["table"])
        else:
            render_bullets(slide, [normalize_item(i) for i in page.get("content", [])])

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    report["output"] = str(output_path)
    return report


EXAMPLE_SPEC = {
    "title": "示例项目\n季度进展汇报",
    "date": "2026年4月1日",
    "output": "example-deck.pptx",
    "outline": [
        {"slide_type": "cover"},
        {
            "slide_type": "content",
            "title": "项目背景",
            "content": [
                "现有流程依赖人工核对，单次处理约 40 分钟",
                "三条业务线各自维护一套口径，数据难以横向比较",
                "本季度目标：统一口径并把处理时间压到 5 分钟内",
            ],
        },
        {
            "slide_type": "content",
            "title": "三个工作方向",
            "content_type": "columns",
            "columns": 3,
            "content": [
                {"title": "统一口径", "items": ["合并三套指标定义", "输出口径说明文档"]},
                {"title": "流程自动化", "items": ["核对环节脚本化", "异常自动告警"]},
                {"title": "结果可视化", "items": ["周报自动生成", "异常按负责人分派"]},
            ],
        },
        {
            "slide_type": "content",
            "title": "关键指标",
            "content_type": "table",
            "table": {
                "headers": ["指标", "上季度", "本季度"],
                "rows": [
                    ["单次处理时长", "40 分钟", "4 分钟"],
                    ["人工介入次数", "每周 12 次", "每周 2 次"],
                    ["口径不一致告警", "无统计", "每周 3 起，均已闭环"],
                ],
            },
        },
        {
            "slide_type": "content",
            "title": "下季度计划",
            "content_type": "mixed",
            "content": ["把方案推广到剩余两条业务线"],
            "table": {
                "headers": ["阶段", "时间", "交付物"],
                "rows": [
                    ["接入改造", "第 1-3 周", "两条线的数据接入"],
                    ["并行验证", "第 4-6 周", "新旧口径对账报告"],
                    ["切换上线", "第 7 周", "旧流程下线"],
                ],
            },
        },
    ],
}


def main() -> int:
    parser = argparse.ArgumentParser(description="Generate a deck from a JSON spec")
    parser.add_argument("spec", nargs="?", help="path to the JSON spec")
    parser.add_argument("--output", help="output .pptx path (overrides spec.output)")
    parser.add_argument("--profile", help="template profile JSON (see profiles/)")
    parser.add_argument("--example", action="store_true", help="print an example spec and exit")
    parser.add_argument("--no-verify", action="store_true", help="skip the built-in layout check")
    args = parser.parse_args()

    if args.example:
        print(json.dumps(EXAMPLE_SPEC, indent=2, ensure_ascii=False))
        return 0
    if not args.spec:
        parser.error("spec is required (or use --example)")

    spec = json.loads(Path(args.spec).read_text(encoding="utf-8"))
    output = args.output or spec.get("output") or "output.pptx"

    try:
        report = generate(spec, output, profile=args.profile)
    except deckprofile.ProfileError as exc:
        print(f"PROFILE ERROR: {exc}", file=sys.stderr)
        return 3
    except SpecError as exc:
        print(f"SPEC ERROR: {exc}", file=sys.stderr)
        return 2

    print(f"wrote {report['output']} ({report['slides']} slides)")
    if report["pages_added"]:
        print(f"  {report['pages_added']} continuation slide(s) added to avoid overflow")
    for note in report["notes"]:
        print(f"  note: {note}")

    if args.no_verify:
        return 0

    import verify_pptx

    findings = verify_pptx.check_file(report["output"])
    verify_pptx.print_report(findings)
    return 1 if any(f["level"] == "error" for f in findings) else 0


if __name__ == "__main__":
    sys.exit(main())
