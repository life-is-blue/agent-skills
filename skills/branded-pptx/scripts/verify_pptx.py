#!/usr/bin/env python3
"""Check a generated deck for the failures python-pptx will not report.

python-pptx happily writes text that spills off the slide, tables taller than
the page, and placeholders that render as "click to add text" prompts. This
script measures the saved file and reports those cases, so the agent has a
feedback loop instead of assuming success.

    python3 verify_pptx.py deck.pptx            # geometry + text fit report
    python3 verify_pptx.py deck.pptx --render out/   # also rasterise to PNG

Exit code is 1 when any error-level finding is present.
"""

from __future__ import annotations

import argparse
import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn
from pptx.util import Emu

sys.path.insert(0, str(Path(__file__).resolve().parent))
import textmetrics as tm  # noqa: E402

DEFAULT_BODY_SIZE_PT = 18.0
DEFAULT_INSET_PT = 7.2  # 0.1 inch
EDGE_TOLERANCE_PT = 4.0
MIN_FILL_RATIO = 0.22
TEXT_PH_TYPES = {
    PP_PLACEHOLDER.BODY,
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
    PP_PLACEHOLDER.SUBTITLE,
    PP_PLACEHOLDER.OBJECT,
}
MAX_FILL_RATIO = 0.95


def _finding(level: str, slide: int, message: str) -> dict:
    return {"level": level, "slide": slide, "message": message}


def _title_placeholder(slide):
    """Same resolution rule as the generator: a text slot at idx 15, else the
    topmost wide text slot in the title band."""
    text_placeholders = [
        shape
        for shape in slide.shapes
        if shape.is_placeholder and shape.placeholder_format.type in TEXT_PH_TYPES
    ]
    for shape in text_placeholders:
        if shape.placeholder_format.idx == 15:
            return shape
    band = [
        shape
        for shape in text_placeholders
        if Emu(shape.top).inches < 1.0 and Emu(shape.width).inches > 4.0
    ]
    return min(band, key=lambda shape: (shape.top, shape.left), default=None)


def _body_pr(text_frame):
    return text_frame._txBody.find(qn("a:bodyPr"))


def _autofit_mode(text_frame) -> str:
    body_pr = _body_pr(text_frame)
    if body_pr is None:
        return "none"
    if body_pr.find(qn("a:spAutoFit")) is not None:
        return "shape"
    if body_pr.find(qn("a:normAutofit")) is not None:
        return "text"
    return "none"


def _wraps(text_frame) -> bool:
    body_pr = _body_pr(text_frame)
    return body_pr is None or body_pr.get("wrap", "square") != "none"


def _paragraph_specs(text_frame) -> list[dict]:
    specs = []
    for paragraph in text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs)
        if not text:
            continue
        sizes = [run.font.size.pt for run in paragraph.runs if run.font.size is not None]
        size = max(sizes) if sizes else DEFAULT_BODY_SIZE_PT
        line_spacing = paragraph.line_spacing if isinstance(paragraph.line_spacing, float) else 1.0
        specs.append(
            {
                "text": text,
                "size_pt": size,
                "line_spacing": line_spacing,
                "space_before_pt": paragraph.space_before.pt if paragraph.space_before else 0.0,
                "space_after_pt": paragraph.space_after.pt if paragraph.space_after else 0.0,
            }
        )
    return specs


def _text_extent(shape) -> tuple[float, float]:
    """Estimated (width_pt, height_pt) of the text inside a shape."""
    specs = _paragraph_specs(shape.text_frame)
    if not specs:
        return 0.0, 0.0
    if _wraps(shape.text_frame):
        width_pt = max(Emu(shape.width).pt - 2 * DEFAULT_INSET_PT, 1.0)
        return width_pt, tm.block_height_pt(specs, width_pt)
    width = max(tm.text_width_pt(spec["text"], spec["size_pt"]) for spec in specs)
    height = sum(
        spec["size_pt"] * tm.LINE_BOX_RATIO * spec["line_spacing"]
        + spec["space_before_pt"]
        + spec["space_after_pt"]
        for spec in specs
    )
    return width, height


def check_file(path: str | Path) -> list[dict]:
    prs = Presentation(str(path))
    slide_width_pt = Emu(prs.slide_width).pt
    slide_height_pt = Emu(prs.slide_height).pt
    findings: list[dict] = []

    for index, slide in enumerate(prs.slides):
        used_top, used_bottom = None, None

        for shape in slide.shapes:
            name = shape.name
            left_pt = Emu(shape.left).pt
            top_pt = Emu(shape.top).pt
            right_pt = left_pt + Emu(shape.width).pt
            bottom_pt = top_pt + Emu(shape.height).pt

            if shape.has_table:
                table_height = sum(Emu(row.height).pt for row in shape.table.rows)
                bottom_pt = top_pt + table_height
                if bottom_pt > slide_height_pt + EDGE_TOLERANCE_PT:
                    findings.append(
                        _finding(
                            "error",
                            index,
                            f"table {name!r} ends {bottom_pt - slide_height_pt:.0f}pt below the "
                            f"slide edge ({len(shape.table.rows)} rows) — paginate it",
                        )
                    )
                for row_index, row in enumerate(shape.table.rows):
                    for col_index, cell in enumerate(row.cells):
                        specs = _paragraph_specs(cell.text_frame)
                        if not specs:
                            continue
                        width_pt = Emu(shape.table.columns[col_index].width).pt - 2 * DEFAULT_INSET_PT
                        needed = tm.block_height_pt(specs, max(width_pt, 1.0))
                        if needed > Emu(row.height).pt * 1.6:
                            findings.append(
                                _finding(
                                    "warning",
                                    index,
                                    f"table cell (r{row_index},c{col_index}) needs about "
                                    f"{needed:.0f}pt in a {Emu(row.height).pt:.0f}pt row — "
                                    "shorten the text or widen the column",
                                )
                            )

            elif shape.has_text_frame and shape.text_frame.text.strip():
                text_width, text_height = _text_extent(shape)
                if not _wraps(shape.text_frame):
                    # A wrap="none" box is grown by the renderer around its own
                    # centre, so overlong text spills off *both* sides — that is
                    # how a cover title ends up clipped at the left edge.
                    centre = left_pt + Emu(shape.width).pt / 2
                    grown = text_width + 2 * DEFAULT_INSET_PT
                    left_edge = min(left_pt, centre - grown / 2)
                    right_edge = max(left_pt + grown, centre + grown / 2)
                    if left_edge < -EDGE_TOLERANCE_PT or right_edge > slide_width_pt + EDGE_TOLERANCE_PT:
                        findings.append(
                            _finding(
                                "error",
                                index,
                                f"non-wrapping text {name!r} grows to "
                                f"{left_edge:.0f}..{right_edge:.0f}pt on a "
                                f"0..{slide_width_pt:.0f}pt slide — it will be clipped; "
                                "give the box explicit width and wrap='square'",
                            )
                        )
                elif _autofit_mode(shape.text_frame) == "shape":
                    bottom_pt = top_pt + text_height
                else:
                    if text_height > Emu(shape.height).pt * 1.05:
                        findings.append(
                            _finding(
                                "error",
                                index,
                                f"text {name!r} needs about {text_height:.0f}pt in a "
                                f"{Emu(shape.height).pt:.0f}pt box — it will overflow",
                            )
                        )
                    bottom_pt = top_pt + min(text_height, Emu(shape.height).pt)

            if (
                shape.is_placeholder
                and shape.placeholder_format.type in (PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.CHART)
                and shape.has_text_frame
                and shape.text_frame.text.strip()
            ):
                findings.append(
                    _finding(
                        "error",
                        index,
                        f"text was written into picture/chart placeholder {name!r} "
                        f"(idx={shape.placeholder_format.idx}) — it will not render as text",
                    )
                )

            elif shape.is_placeholder and shape.has_text_frame and not shape.text_frame.text.strip():
                findings.append(
                    _finding(
                        "warning",
                        index,
                        f"empty placeholder {name!r} will render a 'click to add' prompt — "
                        "fill it or remove it",
                    )
                )

            if bottom_pt > slide_height_pt + EDGE_TOLERANCE_PT or right_pt > slide_width_pt + EDGE_TOLERANCE_PT:
                if not shape.has_table:
                    findings.append(
                        _finding(
                            "error",
                            index,
                            f"shape {name!r} extends past the slide "
                            f"(right={right_pt:.0f}/{slide_width_pt:.0f}pt, "
                            f"bottom={bottom_pt:.0f}/{slide_height_pt:.0f}pt)",
                        )
                    )
            if left_pt < -EDGE_TOLERANCE_PT and Emu(shape.width).pt < slide_width_pt * 0.9:
                findings.append(
                    _finding("warning", index, f"shape {name!r} starts left of the slide edge")
                )

            # Density is only meaningful for content this script drew itself.
            # Placeholder-driven layouts were positioned by a designer.
            if not shape.is_placeholder and (
                shape.has_table or (shape.has_text_frame and shape.text_frame.text.strip())
            ):
                used_top = top_pt if used_top is None else min(used_top, top_pt)
                used_bottom = bottom_pt if used_bottom is None else max(used_bottom, bottom_pt)

        title = _title_placeholder(slide)
        if title is not None:
            has_placeholder_body = any(
                shape.is_placeholder
                and shape.element is not title.element
                and (
                    shape.element.tag.endswith("}pic")
                    or shape.has_table
                    or (shape.has_text_frame and shape.text_frame.text.strip())
                )
                for shape in slide.shapes
            )
            if used_top is None:
                if not has_placeholder_body:
                    findings.append(
                        _finding("error", index, "content slide has a title but no body content")
                    )
            elif has_placeholder_body:
                pass  # mixed placeholder layout; density is the layout's business
            else:
                body_top = Emu(title.top + title.height).pt + 21.6
                available = slide_height_pt - 18.0 - body_top
                ratio = (used_bottom - used_top) / available if available > 0 else 0.0
                if ratio < MIN_FILL_RATIO:
                    findings.append(
                        _finding(
                            "warning",
                            index,
                            f"body area only {ratio:.0%} filled — merge with another slide or "
                            "add a table / picture",
                        )
                    )
                elif ratio > MAX_FILL_RATIO:
                    findings.append(
                        _finding("warning", index, f"body area {ratio:.0%} filled — leave more whitespace")
                    )

    return findings


def render(path: str | Path, out_dir: str | Path) -> list[Path]:
    """Rasterise the deck so the result can actually be looked at."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise RuntimeError("LibreOffice (soffice) is not installed; cannot render")
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(path)],
        check=True,
        capture_output=True,
    )
    pdf = out_dir / (Path(path).stem + ".pdf")
    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        return [pdf]
    subprocess.run(
        [pdftoppm, "-png", "-r", "96", str(pdf), str(out_dir / "slide")],
        check=True,
        capture_output=True,
    )
    return sorted(out_dir.glob("slide-*.png"))


def print_report(findings: list[dict]) -> None:
    errors = [f for f in findings if f["level"] == "error"]
    warnings = [f for f in findings if f["level"] == "warning"]
    for finding in errors + warnings:
        print(f"  {finding['level'].upper():7} slide {finding['slide']}: {finding['message']}")
    if not findings:
        print("  layout check: no issues found")
    else:
        print(f"  layout check: {len(errors)} error(s), {len(warnings)} warning(s)")


def main() -> int:
    parser = argparse.ArgumentParser(description="Verify a generated .pptx")
    parser.add_argument("pptx")
    parser.add_argument("--render", metavar="DIR", help="also rasterise slides to PNG in DIR")
    args = parser.parse_args()

    findings = check_file(args.pptx)
    print_report(findings)

    if args.render:
        try:
            images = render(args.pptx, args.render)
        except (RuntimeError, subprocess.CalledProcessError) as exc:
            print(f"  render failed: {exc}")
            return 1 if any(f["level"] == "error" for f in findings) else 0
        print(f"  rendered {len(images)} file(s) into {args.render}")

    return 1 if any(f["level"] == "error" for f in findings) else 0


if __name__ == "__main__":
    sys.exit(main())
