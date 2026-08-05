#!/usr/bin/env python3
"""Introspect a template: layouts, placeholders, theme, sample slides, profile.

Run this when adopting a template, when one is replaced, or when a shape name
stops resolving. The generator holds no template facts of its own — they all
come from a profile, and this script is how you derive one.

    python3 inspect_template.py --template deck.pptx --profile   # draft profile
    python3 inspect_template.py --layouts            # layout catalogue
    python3 inspect_template.py --layout NAME        # one layout in detail
    python3 inspect_template.py --slides             # sample slides in the file
    python3 inspect_template.py --theme              # theme colours and fonts
    python3 inspect_template.py --json               # everything, machine readable

--profile detection is a heuristic. Read the result before trusting it; the
fields are documented in deckprofile.py.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

SKILL_DIR = Path(__file__).resolve().parent.parent
NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
TEXT_PH_TYPES = {"BODY", "TITLE", "CENTER_TITLE", "SUBTITLE", "OBJECT"}


def inches(value) -> float:
    return round(Emu(int(value)).inches, 2)


def placeholder_rows(container) -> list[dict]:
    rows = []
    for shape in container.placeholders:
        fmt = shape.placeholder_format
        rows.append(
            {
                "idx": fmt.idx,
                "type": str(fmt.type).split(" ")[0],
                "name": shape.name,
                "left": inches(shape.left),
                "top": inches(shape.top),
                "width": inches(shape.width),
                "height": inches(shape.height),
            }
        )
    return sorted(rows, key=lambda row: (row["top"], row["left"]))


def collect_layouts(prs: Presentation) -> dict[str, dict]:
    layouts: dict[str, dict] = {}
    for master_index, master in enumerate(prs.slide_masters):
        for layout in master.slide_layouts:
            if layout.name in layouts:
                continue
            rows = placeholder_rows(layout)
            layouts[layout.name] = {
                "name": layout.name,
                "master": master_index,
                "pictures": sum(1 for row in rows if row["type"] == "PICTURE"),
                "charts": sum(1 for row in rows if row["type"] == "CHART"),
                "bodies": sum(1 for row in rows if row["type"] == "BODY"),
                "placeholders": rows,
            }
    return layouts


def collect_theme(prs: Presentation) -> dict:
    master = prs.slide_masters[0]
    theme_part = master.part.part_related_by(
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
    )
    root = etree.fromstring(theme_part.blob)
    colors = {}
    for node in root.find(".//a:clrScheme", NS):
        tag = etree.QName(node.tag).localname
        child = node[0]
        colors[tag] = child.get("val") or child.get("lastClr")
    fonts = {}
    scheme = root.find(".//a:fontScheme", NS)
    for kind in ("majorFont", "minorFont"):
        node = scheme.find(f"a:{kind}", NS)
        fonts[kind] = {
            "latin": node.find("a:latin", NS).get("typeface"),
            "ea": node.find("a:ea", NS).get("typeface"),
        }
    return {"colors": colors, "fonts": fonts}


def collect_slides(prs: Presentation) -> list[dict]:
    slides = []
    for index, slide in enumerate(prs.slides):
        shapes = []
        for shape in slide.shapes:
            shapes.append(
                {
                    "name": shape.name,
                    "type": str(shape.shape_type).split(" ")[0],
                    "placeholder_idx": shape.placeholder_format.idx if shape.is_placeholder else None,
                    "left": inches(shape.left),
                    "top": inches(shape.top),
                    "width": inches(shape.width),
                    "height": inches(shape.height),
                    "text": shape.text_frame.text[:40] if shape.has_text_frame else "",
                }
            )
        slides.append({"index": index, "layout": slide.slide_layout.name, "shapes": shapes})
    return slides


def detect_profile(prs: Presentation, template_path: Path) -> dict:
    """Derive a draft profile from the template's own structure.

    Heuristics, in order of how much they can be trusted:
      - the cover is the first slide whose layout has no text placeholder;
      - on it, the text box with the largest font is the title, the next is the
        date;
      - the safe right edge is the nearest full-height panel edge right of the
        title, so the title cannot grow over the artwork beside it;
      - the content layout is whichever layout the sample slides use most;
      - the body area starts under that layout's title slot and ends above the
        slide-number placeholder.
    """
    slides = list(prs.slides)
    if not slides:
        raise SystemExit("template has no slides; a cover slide is required")

    cover_index = 0
    for index, slide in enumerate(slides):
        has_text_ph = any(
            shape.is_placeholder
            and str(shape.placeholder_format.type).split(" ")[0] in TEXT_PH_TYPES
            for shape in slide.shapes
        )
        if not has_text_ph:
            cover_index = index
            break
    cover = slides[cover_index]

    def max_font(shape) -> float:
        if not shape.has_text_frame:
            return 0.0
        sizes = [
            run.font.size.pt
            for paragraph in shape.text_frame.paragraphs
            for run in paragraph.runs
            if run.font.size is not None
        ]
        return max(sizes, default=0.0)

    text_boxes = sorted(
        (shape for shape in cover.shapes if shape.has_text_frame and shape.text_frame.text.strip()),
        key=max_font,
        reverse=True,
    )
    if not text_boxes:
        raise SystemExit(f"slide {cover_index} has no text box to use as the cover title")
    title = text_boxes[0]
    date = text_boxes[1] if len(text_boxes) > 1 else None

    # Only full-height artwork counts as a panel edge; small decorations sit
    # beside the title without constraining it.
    title_right = title.left + title.width
    panel_edges = [
        shape.left + shape.width
        for shape in cover.shapes
        if not (shape.has_text_frame and shape.text_frame.text.strip())
        and shape.height > prs.slide_height * 0.5
        and shape.width > prs.slide_width * 0.3
        and shape.left + shape.width > title_right
    ]
    safe_right = min(panel_edges, default=prs.slide_width)

    counts: dict[str, int] = {}
    for index, slide in enumerate(slides):
        if index != cover_index:
            counts[slide.slide_layout.name] = counts.get(slide.slide_layout.name, 0) + 1
    content_layout = max(counts, key=counts.get) if counts else cover.slide_layout.name

    layout = next(
        (
            candidate
            for master in prs.slide_masters
            for candidate in master.slide_layouts
            if candidate.name == content_layout
        ),
        None,
    )
    title_idx = 15
    bottom_in = round(Emu(prs.slide_height).inches - 0.6, 2)
    if layout is not None:
        band = [
            shape
            for shape in layout.placeholders
            if str(shape.placeholder_format.type).split(" ")[0] in TEXT_PH_TYPES
            and inches(shape.top) < 1.0
            and inches(shape.width) > 4.0
        ]
        if band:
            title_idx = min(band, key=lambda shape: (shape.top, shape.left)).placeholder_format.idx
        numbers = [
            shape
            for shape in layout.placeholders
            if str(shape.placeholder_format.type).startswith("SLIDE_NUMBER")
        ]
        if numbers:
            bottom_in = round(min(inches(shape.top) for shape in numbers) - 0.25, 2)

    return {
        "template": template_path.name,
        "cover_slide_index": cover_index,
        "cover_title_shape": title.name,
        "cover_date_shape": date.name if date is not None else "",
        "cover_safe_right_in": inches(safe_right),
        "cover_safe_margin_in": 0.3,
        "content_layout": content_layout,
        "title_placeholder_idx": title_idx,
        "content_margin_x_in": 0.5,
        "content_gap_below_title_in": 0.3,
        "content_bottom_in": bottom_in,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Inspect a PPTX template")
    parser.add_argument("--template", required=True, help="path to the template .pptx")
    parser.add_argument("--layouts", action="store_true")
    parser.add_argument("--layout", metavar="NAME")
    parser.add_argument("--slides", action="store_true")
    parser.add_argument("--theme", action="store_true")
    parser.add_argument("--json", action="store_true")
    parser.add_argument("--profile", action="store_true", help="print a draft profile JSON")
    args = parser.parse_args()

    prs = Presentation(args.template)
    layouts = collect_layouts(prs)

    if args.profile:
        profile = detect_profile(prs, Path(args.template))
        print(json.dumps(profile, indent=2, ensure_ascii=False))
        print(
            "\n# Review every value above, then save it as profiles/<name>.json and "
            "\n# put the template where 'template' points (relative to the profile).",
            file=sys.stderr,
        )
        return 0

    if args.json:
        print(
            json.dumps(
                {
                    "slide_size": [inches(prs.slide_width), inches(prs.slide_height)],
                    "theme": collect_theme(prs),
                    "layouts": layouts,
                    "slides": collect_slides(prs),
                },
                indent=2,
                ensure_ascii=False,
            )
        )
        return 0

    if args.layout:
        layout = layouts.get(args.layout)
        if layout is None:
            print(f"no such layout: {args.layout}", file=sys.stderr)
            return 1
        print(f"{layout['name']} (master {layout['master']})")
        for row in layout["placeholders"]:
            print(
                f"  idx={row['idx']:<3} {row['type']:<13} "
                f"({row['left']}, {row['top']}) {row['width']}x{row['height']}in  {row['name']}"
            )
        return 0

    if args.theme:
        theme = collect_theme(prs)
        for key, value in theme["colors"].items():
            print(f"  {key:<9} #{value}")
        for kind, value in theme["fonts"].items():
            print(f"  {kind:<9} latin={value['latin']} ea={value['ea']}")
        return 0

    if args.slides:
        for slide in collect_slides(prs):
            print(f"slide {slide['index']} — layout {slide['layout']}")
            for shape in slide["shapes"]:
                print(
                    f"  {shape['type']:<12} idx={shape['placeholder_idx']} "
                    f"({shape['left']}, {shape['top']}) {shape['width']}x{shape['height']}in "
                    f"{shape['name']!r} {shape['text']!r}"
                )
        return 0

    print(f"{Path(args.template).name}: {inches(prs.slide_width)} x {inches(prs.slide_height)}in, "
          f"{len(prs.slides)} sample slides, {len(layouts)} unique layouts")
    for layout in sorted(layouts.values(), key=lambda item: item["name"]):
        print(
            f"  {layout['name']:<20} bodies={layout['bodies']:<3} "
            f"pictures={layout['pictures']:<3} charts={layout['charts']}"
        )
    return 0


if __name__ == "__main__":
    sys.exit(main())
